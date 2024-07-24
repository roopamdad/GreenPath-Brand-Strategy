from crewai import Agent, Task
from pptx import Presentation

class DocumentationAgent(Agent):
    @Task()
    def create_presentation(self, output_pptx_path: str):
        # Create a presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"

        # Add content
        content = slide.shapes.add_textbox(left=0, top=1, width=prs.slide_width, height=prs.slide_height)
        text_frame = content.text_frame
        text_frame.text = "Summary of findings and recommendations."

        # Save presentation
        prs.save(output_pptx_path)
